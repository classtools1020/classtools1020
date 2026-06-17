# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

# palette
PEACH=RGBColor(0xFF,0xB7,0x4D); LAV=RGBColor(0x9F,0x86,0xC9); GREEN=RGBColor(0x4C,0xAF,0x50)
RED=RGBColor(0xE5,0x39,0x35); ORANGE=RGBColor(0xFB,0x8C,0x00); INDIGO=RGBColor(0x5C,0x6B,0xC0)
TEAL=RGBColor(0x00,0x89,0x7B); CREAM=RGBColor(0xFF,0xF8,0xE1); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK=RGBColor(0x37,0x37,0x37); CORAL=RGBColor(0xFF,0x70,0x43); PURPLE=RGBColor(0x6A,0x1B,0x9A)
PINK=RGBColor(0xF0,0x62,0x92)

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, x, y, w, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp=1.1):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i,(text,size,bold,color) in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp; p.space_after = Pt(6)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
        f.name = 'Microsoft JhengHei'
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {'typeface': 'Microsoft JhengHei'})
        rPr.append(ea)
    return box

def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def footer(slide, n, label):
    tb(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4),
       [(label, 12, False, RGBColor(0xAA,0xAA,0xAA))], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)
    tb(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
       [(str(n), 12, False, RGBColor(0xAA,0xAA,0xAA))], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

UNIT='必須說出的祕密｜國中特教班社會科'

def s_cover(t, sub, c1, c2):
    s = prs.slides.add_slide(BLANK); bg(s, c1)
    tb(s, Inches(1), Inches(2.2), Inches(11.33), Inches(2), [(t, 54, True, WHITE)])
    tb(s, Inches(1), Inches(4.4), Inches(11.33), Inches(1.2), [(sub, 28, False, WHITE)])
    return s

def s_color(t, sub, c, fg=WHITE):
    s = prs.slides.add_slide(BLANK); bg(s, c)
    tb(s, Inches(1), Inches(2.3), Inches(11.33), Inches(1.8), [(t, 48, True, fg)])
    if sub:
        tb(s, Inches(1.2), Inches(4.2), Inches(10.9), Inches(1.8), [(sub, 30, False, fg)])
    return s

def s_content(t, bullets, accent):
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(1.5), accent)
    tb(s, Inches(0.6), Inches(0.1), Inches(12.1), Inches(1.3), [(t, 40, True, WHITE)])
    runs = [(b, 30, False, DARK) for b in bullets]
    tb(s, Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.3), runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.4)
    return s

def s_question(num, title, scenario, options):
    s = prs.slides.add_slide(BLANK); bg(s, TEAL)
    tb(s, Inches(0.6), Inches(0.4), Inches(12), Inches(1), [(f'❓ 闖關題目 {num}', 30, True, RGBColor(0xB2,0xDF,0xDB))])
    tb(s, Inches(1), Inches(1.5), Inches(11.33), Inches(1.4), [(title, 40, True, WHITE)])
    tb(s, Inches(1.2), Inches(3.0), Inches(10.9), Inches(1.6), [(scenario, 30, False, WHITE)])
    rect(s, Inches(2.5), Inches(5.2), Inches(8.33), Inches(1.3), WHITE)
    tb(s, Inches(2.5), Inches(5.2), Inches(8.33), Inches(1.3), [(options, 30, True, TEAL)])
    return s

def s_answer(num, big, detail, color):
    s = prs.slides.add_slide(BLANK); bg(s, CREAM)
    tb(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.9), [(f'✅ 答案 {num}', 26, True, RGBColor(0xBD,0xA1,0x55))])
    tb(s, Inches(1), Inches(1.8), Inches(11.33), Inches(2), [(big, 54, True, color)])
    tb(s, Inches(1.2), Inches(4.2), Inches(10.9), Inches(2), [(detail, 32, False, DARK)])
    return s

def s_video(t, url, sub=''):
    s = prs.slides.add_slide(BLANK); bg(s, INDIGO)
    tb(s, Inches(1), Inches(1.6), Inches(11.33), Inches(1.5), [('▶ '+t, 40, True, WHITE)])
    if sub:
        tb(s, Inches(1), Inches(3.1), Inches(11.33), Inches(1), [(sub, 28, False, RGBColor(0xC5,0xCA,0xE9))])
    rect(s, Inches(1.5), Inches(4.3), Inches(10.33), Inches(1.1), WHITE)
    tb(s, Inches(1.5), Inches(4.3), Inches(10.33), Inches(1.1), [(url, 22, True, INDIGO)])
    return s

# ---- build 45 slides ----
slides = [
 ('cover','必須說出的祕密 🤫','好祕密 vs 壞祕密・勇敢說出來',PINK,None),
 ('color','你有沒有祕密呢？','每個人都有祕密，但祕密有兩種喔！',LAV),
 ('content','祕密有兩種 🟢🔴',['🟢 好祕密：讓人開心','🔴 壞祕密：讓你害怕、不舒服'],PINK),
 ('color','🟢 好祕密','是驚喜（生日禮物、派對）\n之後會說出來，可以收著',GREEN),
 ('color','🔴 壞祕密','讓你害怕、不舒服\n叫你不能告訴任何人 → 要說出來',RED),
 ('content','危險訊號① ⚠️',['有人說「這是我們的祕密，不能說」'],ORANGE),
 ('content','危險訊號② 🎁',['用禮物哄你保密，要小心'],ORANGE),
 ('content','危險訊號③ 😨',['用威脅叫你保密，更要小心'],ORANGE),
 ('content','身體會報警 💓',['睡不著、肚子痛、怕怕的','那是身體在說：不對勁！'],CORAL),
 ('video','看影片：我說可以才可以','https://www.youtube.com/watch?v=lXCrxuEl-Kw','身體自主權・易讀版動畫'),
 ('content','影片重點',['我說可以，才可以！','誰可以決定？自己！'],TEAL),
 ('color','我的身體我決定','我的身體、我的祕密，我做主',PURPLE),
 ('question','①','生日驚喜先別說？','朋友要給你驚喜，叫你先別說。','🟢 好祕密？　🔴 壞祕密？'),
 ('answer','①','🟢 好祕密！','之後會公開，可以保守。',GREEN),
 ('content','小結：好祕密會公開',['好祕密可以放心收著 🟢'],GREEN),
 ('color','壞祕密大搜查 🔍','接下來六個情境，一起來判斷！',INDIGO),
 ('question','②','有人碰你叫你別說？','有人碰你的身體，叫你不能告訴爸媽。','🟢 好祕密？　🔴 壞祕密？'),
 ('answer','②','🔴 壞祕密！','一定要說出來。',RED),
 ('question','③','答應了讓你害怕的祕密？','你答應幫他保守，但心裡很害怕。','🟢 好祕密？　🔴 壞祕密？'),
 ('answer','③','🔴 還是要說！','就算你答應了，也要說出來。',RED),
 ('question','④','睡不著、肚子痛的祕密？','這個祕密讓你身體很不舒服。','🟢 好祕密？　🔴 壞祕密？'),
 ('answer','④','🔴 壞祕密！','身體在報警，要說出來。',RED),
 ('question','⑤','網路叫你傳照片別說？','網路上有人叫你傳照片、別告訴爸媽。','🟢 好祕密？　🔴 壞祕密？'),
 ('answer','⑤','🔴 網路壞祕密！','立刻告訴大人。',RED),
 ('question','⑥','被威脅別說出去？','他威脅你「說出去你會有麻煩」。','🟢 好祕密？　🔴 壞祕密？'),
 ('answer','⑥','🔴 更要說！','威脅你保密，更要說出來。',RED),
 ('video','〔A組〕新聞識讀《不能說的祕密》','https://www.youtube.com/watch?v=efvmDxR7240','媒體識讀練習'),
 ('content','〔A組〕識讀四問（上）',['① 這是事實還是意見？','② 有沒有保護當事人？'],TEAL),
 ('content','〔A組〕識讀四問（下）',['③ 標題會誇大嗎？','④ 我能做什麼？'],TEAL),
 ('content','小結：壞祕密都要說',['每一種壞祕密，都要說出來 🔴'],RED),
 ('color','勇敢說出來 🗣️','壞祕密，我可以找誰幫忙？',INDIGO),
 ('content','我可以找誰說？',['找「信任的大人」'],PURPLE),
 ('color','信任的大人① 爸爸媽媽','最親近、最信任的人',GREEN),
 ('color','信任的大人② 老師、輔導老師','在學校可以找的人',TEAL),
 ('color','信任的大人③ 阿公阿嬤','家裡可以信任的長輩',LAV),
 ('content','動手畫：我的 3 個信任大人',['寫或畫出我的求助對象','①________ ②________ ③________'],PINK),
 ('content','說不出口怎麼辦？',['可以寫紙條給信任的大人'],CORAL),
 ('color','打 113 📞','24小時兒少保護專線',RED),
 ('content','法律也保護你 ⚖️',['兒少權益保障法（公民連結）','法律站在你這邊！'],INDIGO),
 ('content','三步驟',['認出來 → 說出來 → 找大人'],GREEN),
 ('content','勇敢練習',['練習說：「這個祕密讓我不舒服。」'],TEAL),
 ('content','角色扮演',['兩人一組，跟信任的大人說說看'],PURPLE),
 ('content','手作：勇敢說出來求助卡',['寫上信任大人和 113，帶在身邊'],PINK),
 ('color','說出來，是保護自己 💪','說出來不是告密，是保護自己',CORAL),
 ('cover','我會勇敢說出來！','好祕密🟢收，壞祕密🔴一定說',PURPLE,None),
]

for i, sl in enumerate(slides, 1):
    typ = sl[0]
    if typ=='cover':
        s = s_cover(sl[1], sl[2], sl[3], None)
    elif typ=='color':
        s = s_color(sl[1], sl[2], sl[3])
    elif typ=='content':
        s = s_content(sl[1], sl[2], sl[3])
    elif typ=='question':
        s = s_question(sl[1], sl[2], sl[3], sl[4])
    elif typ=='answer':
        s = s_answer(sl[1], sl[2], sl[3], sl[4])
    elif typ=='video':
        s = s_video(sl[1], sl[2], sl[3])
    footer(s, i, UNIT)

out='/home/user/classtools1020/必須說出的祕密_45張簡報.pptx'
prs.save(out); print('saved', out, 'slides:', len(slides))
